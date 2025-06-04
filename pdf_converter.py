#!/usr/bin/env python3
"""
PDF Converter Desktop App for Mac
Converts files in a directory and subdirectories to PDF format
"""

import os
import sys
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pathlib import Path
import threading
from datetime import datetime

# Import conversion libraries
try:
    from PIL import Image
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    import markdown
except ImportError as e:
    print(f"Error importing required libraries: {e}")
    print("Please install requirements: pip install -r requirements.txt")
    sys.exit(1)


class PDFConverter:
    def __init__(self):
        self.supported_formats = {
            'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'],
            'documents': ['.docx', '.txt', '.md'],
            'spreadsheets': ['.xlsx', '.xls'],
            'presentations': ['.pptx'],
            'pdf': ['.pdf']  # For combining/copying existing PDFs
        }
        
    def convert_image_to_pdf(self, input_path, output_path):
        """Convert image files to PDF"""
        try:
            with Image.open(input_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                img.save(output_path, "PDF", resolution=100.0)
            return True
        except Exception as e:
            print(f"Error converting image {input_path}: {e}")
            return False
    
    def convert_docx_to_pdf(self, input_path, output_path):
        """Convert DOCX files to PDF"""
        try:
            doc = Document(input_path)
            
            # Create PDF with reportlab
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            if story:
                pdf_doc.build(story)
            else:
                # Create empty PDF if no content
                c = canvas.Canvas(output_path, pagesize=A4)
                c.drawString(100, 750, f"Converted from: {os.path.basename(input_path)}")
                c.drawString(100, 730, "No readable content found")
                c.save()
            
            return True
        except Exception as e:
            print(f"Error converting DOCX {input_path}: {e}")
            return False
    
    def convert_txt_to_pdf(self, input_path, output_path):
        """Convert text files to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Split content into lines and create paragraphs
            lines = content.split('\n')
            for line in lines:
                if line.strip():
                    p = Paragraph(line, styles['Normal'])
                    story.append(p)
                else:
                    story.append(Spacer(1, 12))
            
            if story:
                pdf_doc.build(story)
            else:
                # Create empty PDF if no content
                c = canvas.Canvas(output_path, pagesize=A4)
                c.drawString(100, 750, f"Converted from: {os.path.basename(input_path)}")
                c.save()
            
            return True
        except Exception as e:
            print(f"Error converting TXT {input_path}: {e}")
            return False
    
    def convert_md_to_pdf(self, input_path, output_path):
        """Convert Markdown files to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to HTML then to PDF
            html_content = markdown.markdown(md_content)
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Simple conversion - just treat as text for now
            p = Paragraph(html_content, styles['Normal'])
            story.append(p)
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting MD {input_path}: {e}")
            return False
    
    def convert_xlsx_to_pdf(self, input_path, output_path):
        """Convert Excel files to PDF"""
        try:
            workbook = load_workbook(input_path)
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            y_position = height - 50
            
            c.drawString(50, y_position, f"Excel File: {os.path.basename(input_path)}")
            y_position -= 30
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                c.drawString(50, y_position, f"Sheet: {sheet_name}")
                y_position -= 20
                
                for row in sheet.iter_rows(max_row=50, max_col=10, values_only=True):
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        c.drawString(50, y_position, row_text[:80])  # Limit text length
                        y_position -= 15
                
                y_position -= 20
            
            c.save()
            return True
        except Exception as e:
            print(f"Error converting XLSX {input_path}: {e}")
            return False
    
    def convert_pptx_to_pdf(self, input_path, output_path):
        """Convert PowerPoint files to PDF"""
        try:
            presentation = Presentation(input_path)
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.drawString(50, height - 50, f"PowerPoint: {os.path.basename(input_path)}")
            
            slide_num = 1
            for slide in presentation.slides:
                c.showPage()
                c.drawString(50, height - 50, f"Slide {slide_num}")
                y_position = height - 80
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_lines = shape.text.split('\n')
                        for line in text_lines:
                            if y_position < 50:
                                c.showPage()
                                y_position = height - 50
                            c.drawString(50, y_position, line[:80])
                            y_position -= 20
                
                slide_num += 1
            
            c.save()
            return True
        except Exception as e:
            print(f"Error converting PPTX {input_path}: {e}")
            return False
    
    def get_file_type(self, file_path):
        """Determine the file type category"""
        extension = Path(file_path).suffix.lower()
        
        for category, extensions in self.supported_formats.items():
            if extension in extensions:
                return category
        return None
    
    def convert_file(self, input_path, output_path):
        """Convert a single file to PDF"""
        file_type = self.get_file_type(input_path)
        
        if not file_type:
            return False
        
        try:
            if file_type == 'images':
                return self.convert_image_to_pdf(input_path, output_path)
            elif file_type == 'documents':
                ext = Path(input_path).suffix.lower()
                if ext == '.docx':
                    return self.convert_docx_to_pdf(input_path, output_path)
                elif ext == '.txt':
                    return self.convert_txt_to_pdf(input_path, output_path)
                elif ext == '.md':
                    return self.convert_md_to_pdf(input_path, output_path)
            elif file_type == 'spreadsheets':
                return self.convert_xlsx_to_pdf(input_path, output_path)
            elif file_type == 'presentations':
                return self.convert_pptx_to_pdf(input_path, output_path)
            elif file_type == 'pdf':
                # Copy existing PDF
                import shutil
                shutil.copy2(input_path, output_path)
                return True
            
            return False
        except Exception as e:
            print(f"Error converting {input_path}: {e}")
            return False


class PDFConverterGUI:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("PDF Converter for Mac")
        self.root.geometry("600x550")
        self.root.resizable(True, True)
        
        # Initialize converter
        self.converter = PDFConverter()
        
        # Variables
        self.selected_directory = tk.StringVar()
        self.conversion_progress = tk.StringVar(value="Ready to convert files")
        self.maintain_structure = tk.BooleanVar(value=True)  # New variable for directory structure
        
        self.setup_gui()
        
    def setup_gui(self):
        """Setup the GUI components"""
        # Configure root window
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        
        # Main frame
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # Configure main frame grid weights
        main_frame.columnconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        main_frame.columnconfigure(2, weight=1)
        main_frame.rowconfigure(7, weight=1)  # Make results area expandable
        
        # Title
        title_label = ttk.Label(main_frame, text="PDF Converter", 
                               font=("Arial", 20, "bold"))
        title_label.grid(row=0, column=0, columnspan=3, pady=(0, 20), sticky=tk.W+tk.E)
        
        # Directory selection
        dir_label = ttk.Label(main_frame, text="Select Directory:")
        dir_label.grid(row=1, column=0, columnspan=3, sticky=tk.W, pady=(0, 5))
        
        directory_frame = ttk.Frame(main_frame)
        directory_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        directory_frame.columnconfigure(0, weight=1)
        
        self.directory_entry = ttk.Entry(directory_frame, textvariable=self.selected_directory,
                                        font=("Arial", 12))
        self.directory_entry.grid(row=0, column=0, sticky=(tk.W, tk.E), padx=(0, 10))
        
        browse_button = ttk.Button(directory_frame, text="Browse", 
                                  command=self.browse_directory)
        browse_button.grid(row=0, column=1, sticky=tk.E)
        
        # Output structure options
        output_frame = ttk.LabelFrame(main_frame, text="Output Structure", padding="10")
        output_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        output_frame.columnconfigure(0, weight=1)
        
        structure_radio1 = ttk.Radiobutton(output_frame, text="Maintain directory structure", 
                                          variable=self.maintain_structure, value=True)
        structure_radio1.grid(row=0, column=0, sticky=tk.W, padx=10, pady=2)
        
        structure_radio2 = ttk.Radiobutton(output_frame, text="Flat list (all PDFs in one folder)", 
                                          variable=self.maintain_structure, value=False)
        structure_radio2.grid(row=1, column=0, sticky=tk.W, padx=10, pady=2)
        
        # Supported formats info
        info_frame = ttk.LabelFrame(main_frame, text="Supported Formats", padding="10")
        info_frame.grid(row=4, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        info_frame.columnconfigure(0, weight=1)
        
        formats_text = """• Images: JPG, PNG, GIF, BMP, TIFF, WebP
• Documents: DOCX, TXT, Markdown
• Spreadsheets: XLSX, XLS
• Presentations: PPTX
• PDFs: Existing PDFs will be copied"""
        
        formats_label = ttk.Label(info_frame, text=formats_text, font=("Arial", 10))
        formats_label.grid(row=0, column=0, sticky=tk.W)
        
        # Progress section
        progress_frame = ttk.LabelFrame(main_frame, text="Progress", padding="10")
        progress_frame.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        progress_frame.columnconfigure(0, weight=1)
        
        self.progress_bar = ttk.Progressbar(progress_frame, mode='determinate')
        self.progress_bar.grid(row=0, column=0, sticky=(tk.W, tk.E), pady=(0, 5))
        
        self.progress_label = ttk.Label(progress_frame, textvariable=self.conversion_progress,
                                       font=("Arial", 10))
        self.progress_label.grid(row=1, column=0, sticky=tk.W)
        
        # Convert button
        self.convert_button = ttk.Button(main_frame, text="Convert to PDF", 
                                        command=self.start_conversion)
        self.convert_button.grid(row=6, column=0, columnspan=3, pady=15)
        
        # Results text area
        results_frame = ttk.LabelFrame(main_frame, text="Conversion Results", padding="10")
        results_frame.grid(row=7, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=(0, 0))
        results_frame.columnconfigure(0, weight=1)
        results_frame.rowconfigure(0, weight=1)
        
        # Text widget with scrollbar
        text_scrollbar_frame = ttk.Frame(results_frame)
        text_scrollbar_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        text_scrollbar_frame.columnconfigure(0, weight=1)
        text_scrollbar_frame.rowconfigure(0, weight=1)
        
        self.results_text = tk.Text(text_scrollbar_frame, height=8, font=("Monaco", 10), 
                                   wrap=tk.WORD)
        scrollbar = ttk.Scrollbar(text_scrollbar_frame, orient="vertical", 
                                 command=self.results_text.yview)
        self.results_text.configure(yscrollcommand=scrollbar.set)
        
        self.results_text.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        scrollbar.grid(row=0, column=1, sticky=(tk.N, tk.S))
        
        # Initial message in results
        self.results_text.insert(tk.END, "Welcome to PDF Converter!\n")
        self.results_text.insert(tk.END, "Select a directory and click 'Convert to PDF' to begin.\n")
        self.results_text.config(state=tk.DISABLED)
    
    def browse_directory(self):
        """Open directory selection dialog"""
        directory = filedialog.askdirectory(title="Select Directory to Convert")
        if directory:
            self.selected_directory.set(directory)
    
    def log_message(self, message):
        """Add message to results text area"""
        self.results_text.config(state=tk.NORMAL)
        self.results_text.insert(tk.END, f"{datetime.now().strftime('%H:%M:%S')} - {message}\n")
        self.results_text.see(tk.END)
        self.results_text.config(state=tk.DISABLED)
        self.root.update()
    
    def start_conversion(self):
        """Start the conversion process in a separate thread"""
        directory = self.selected_directory.get().strip()
        
        if not directory:
            messagebox.showerror("Error", "Please select a directory first")
            return
        
        if not os.path.exists(directory):
            messagebox.showerror("Error", "Selected directory does not exist")
            return
        
        # Clear results
        self.results_text.config(state=tk.NORMAL)
        self.results_text.delete(1.0, tk.END)
        self.results_text.config(state=tk.DISABLED)
        
        # Disable convert button
        self.convert_button.config(state='disabled')
        
        # Start conversion in separate thread
        thread = threading.Thread(target=self.convert_files, args=(directory,))
        thread.daemon = True
        thread.start()
    
    def generate_unique_filename(self, output_dir, base_name, extension=".pdf"):
        """Generate a unique filename to avoid conflicts in flat structure mode"""
        output_path = output_dir / f"{base_name}{extension}"
        counter = 1
        
        while output_path.exists():
            output_path = output_dir / f"{base_name}_{counter}{extension}"
            counter += 1
        
        return output_path
    
    def convert_files(self, directory):
        """Convert all files in directory and subdirectories"""
        try:
            directory_path = Path(directory)
            pdf_output_dir = directory_path / "pdf"
            
            # Create PDF output directory
            pdf_output_dir.mkdir(exist_ok=True)
            
            structure_mode = "directory structure" if self.maintain_structure.get() else "flat list"
            self.log_message(f"Starting conversion for directory: {directory}")
            self.log_message(f"Output mode: {structure_mode}")
            self.log_message(f"PDF output directory: {pdf_output_dir}")
            
            # Get all files to convert
            all_files = []
            for root, dirs, files in os.walk(directory):
                for file in files:
                    file_path = Path(root) / file
                    if self.converter.get_file_type(file_path):
                        all_files.append(file_path)
            
            if not all_files:
                self.log_message("No supported files found in the directory")
                self.conversion_progress.set("No files to convert")
                self.convert_button.config(state='normal')
                return
            
            self.log_message(f"Found {len(all_files)} files to convert")
            
            # Setup progress bar
            self.progress_bar.config(maximum=len(all_files))
            
            successful_conversions = 0
            failed_conversions = 0
            
            for i, file_path in enumerate(all_files):
                try:
                    # Create relative path structure in PDF directory
                    relative_path = file_path.relative_to(directory_path)
                    pdf_filename = relative_path.stem + ".pdf"
                    
                    if self.maintain_structure.get():
                        # Maintain directory structure
                        output_subdir = pdf_output_dir / relative_path.parent
                        output_subdir.mkdir(parents=True, exist_ok=True)
                        output_path = output_subdir / pdf_filename
                    else:
                        # Flat structure - all files in root pdf directory
                        # Generate unique filename to avoid conflicts
                        base_name = relative_path.stem
                        # If file is in subdirectory, include parent directory name to make it unique
                        if relative_path.parent != Path('.'):
                            parent_name = str(relative_path.parent).replace('/', '_').replace('\\', '_')
                            base_name = f"{parent_name}_{base_name}"
                        
                        output_path = self.generate_unique_filename(pdf_output_dir, base_name)
                    
                    self.conversion_progress.set(f"Converting: {file_path.name}")
                    
                    # Convert file
                    if self.converter.convert_file(str(file_path), str(output_path)):
                        successful_conversions += 1
                        if self.maintain_structure.get():
                            self.log_message(f"✓ Converted: {relative_path}")
                        else:
                            self.log_message(f"✓ Converted: {relative_path} → {output_path.name}")
                    else:
                        failed_conversions += 1
                        self.log_message(f"✗ Failed: {relative_path}")
                    
                    # Update progress
                    self.progress_bar.config(value=i + 1)
                    self.root.update()
                    
                except Exception as e:
                    failed_conversions += 1
                    self.log_message(f"✗ Error converting {file_path.name}: {str(e)}")
            
            # Final summary
            self.log_message(f"\nConversion completed!")
            self.log_message(f"Successfully converted: {successful_conversions} files")
            self.log_message(f"Failed conversions: {failed_conversions} files")
            self.log_message(f"PDFs saved to: {pdf_output_dir}")
            
            self.conversion_progress.set(f"Completed: {successful_conversions} converted, {failed_conversions} failed")
            
            # Show completion message
            structure_info = "maintaining directory structure" if self.maintain_structure.get() else "in flat structure"
            messagebox.showinfo("Conversion Complete", 
                              f"Conversion completed!\n\n"
                              f"Successfully converted: {successful_conversions} files\n"
                              f"Failed conversions: {failed_conversions} files\n\n"
                              f"PDFs saved to: {pdf_output_dir}\n"
                              f"Structure: {structure_info}")
            
        except Exception as e:
            self.log_message(f"Error during conversion: {str(e)}")
            messagebox.showerror("Error", f"An error occurred during conversion:\n{str(e)}")
        
        finally:
            # Re-enable convert button
            self.convert_button.config(state='normal')
    
    def run(self):
        """Run the GUI application"""
        self.root.mainloop()


def main():
    """Main function to run the PDF Converter app"""
    try:
        app = PDFConverterGUI()
        app.run()
    except Exception as e:
        print(f"Error starting application: {e}")
        messagebox.showerror("Error", f"Failed to start application:\n{str(e)}")


if __name__ == "__main__":
    main() 