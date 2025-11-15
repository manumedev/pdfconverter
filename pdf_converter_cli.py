#!/usr/bin/env python3
"""
PDF Converter Command Line Tool for Mac
Converts files in a directory and subdirectories to PDF format
"""

import os
import sys
import argparse
from pathlib import Path
from datetime import datetime

# Import conversion libraries
try:
    from PIL import Image
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    import markdown
except ImportError as e:
    print(f"Error importing required libraries: {e}")
    print("Please install requirements: pip install -r requirements.txt")
    sys.exit(1)


class PDFConverter:
    def __init__(self):
        self.supported_formats = {
            'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'],
            'documents': ['.docx', '.txt', '.md'],
            'spreadsheets': ['.xlsx', '.xls'],
            'presentations': ['.pptx'],
            'pdf': ['.pdf']  # For combining/copying existing PDFs
        }
        
    def convert_image_to_pdf(self, input_path, output_path):
        """Convert image files to PDF"""
        try:
            with Image.open(input_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                img.save(output_path, "PDF", resolution=100.0)
            return True
        except Exception as e:
            print(f"Error converting image {input_path}: {e}")
            return False
    
    def convert_docx_to_pdf(self, input_path, output_path):
        """Convert DOCX files to PDF"""
        try:
            doc = Document(input_path)
            
            # Create PDF with reportlab
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            if story:
                pdf_doc.build(story)
            else:
                # Create empty PDF if no content
                c = canvas.Canvas(output_path, pagesize=A4)
                c.drawString(100, 750, f"Converted from: {os.path.basename(input_path)}")
                c.drawString(100, 730, "No readable content found")
                c.save()
            
            return True
        except Exception as e:
            print(f"Error converting DOCX {input_path}: {e}")
            return False
    
    def convert_txt_to_pdf(self, input_path, output_path):
        """Convert text files to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Split content into lines and create paragraphs
            lines = content.split('\n')
            for line in lines:
                if line.strip():
                    p = Paragraph(line, styles['Normal'])
                    story.append(p)
                else:
                    story.append(Spacer(1, 12))
            
            if story:
                pdf_doc.build(story)
            else:
                # Create empty PDF if no content
                c = canvas.Canvas(output_path, pagesize=A4)
                c.drawString(100, 750, f"Converted from: {os.path.basename(input_path)}")
                c.save()
            
            return True
        except Exception as e:
            print(f"Error converting TXT {input_path}: {e}")
            return False
    
    def convert_md_to_pdf(self, input_path, output_path):
        """Convert Markdown files to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to HTML then to PDF
            html_content = markdown.markdown(md_content)
            
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Simple conversion - just treat as text for now
            p = Paragraph(html_content, styles['Normal'])
            story.append(p)
            
            pdf_doc.build(story)
            return True
        except Exception as e:
            print(f"Error converting MD {input_path}: {e}")
            return False
    
    def convert_xlsx_to_pdf(self, input_path, output_path):
        """Convert Excel files to PDF"""
        try:
            workbook = load_workbook(input_path)
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            y_position = height - 50
            
            c.drawString(50, y_position, f"Excel File: {os.path.basename(input_path)}")
            y_position -= 30
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                c.drawString(50, y_position, f"Sheet: {sheet_name}")
                y_position -= 20
                
                for row in sheet.iter_rows(max_row=50, max_col=10, values_only=True):
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        c.drawString(50, y_position, row_text[:80])  # Limit text length
                        y_position -= 15
                
                y_position -= 20
            
            c.save()
            return True
        except Exception as e:
            print(f"Error converting XLSX {input_path}: {e}")
            return False
    
    def convert_pptx_to_pdf(self, input_path, output_path):
        """Convert PowerPoint files to PDF"""
        try:
            presentation = Presentation(input_path)
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            c.drawString(50, height - 50, f"PowerPoint: {os.path.basename(input_path)}")
            
            slide_num = 1
            for slide in presentation.slides:
                c.showPage()
                c.drawString(50, height - 50, f"Slide {slide_num}")
                y_position = height - 80
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_lines = shape.text.split('\n')
                        for line in text_lines:
                            if y_position < 50:
                                c.showPage()
                                y_position = height - 50
                            c.drawString(50, y_position, line[:80])
                            y_position -= 20
                
                slide_num += 1
            
            c.save()
            return True
        except Exception as e:
            print(f"Error converting PPTX {input_path}: {e}")
            return False
    
    def get_file_type(self, file_path):
        """Determine the file type category"""
        extension = Path(file_path).suffix.lower()
        
        for category, extensions in self.supported_formats.items():
            if extension in extensions:
                return category
        return None
    
    def convert_file(self, input_path, output_path):
        """Convert a single file to PDF"""
        file_type = self.get_file_type(input_path)
        
        if not file_type:
            return False
        
        try:
            if file_type == 'images':
                return self.convert_image_to_pdf(input_path, output_path)
            elif file_type == 'documents':
                ext = Path(input_path).suffix.lower()
                if ext == '.docx':
                    return self.convert_docx_to_pdf(input_path, output_path)
                elif ext == '.txt':
                    return self.convert_txt_to_pdf(input_path, output_path)
                elif ext == '.md':
                    return self.convert_md_to_pdf(input_path, output_path)
            elif file_type == 'spreadsheets':
                return self.convert_xlsx_to_pdf(input_path, output_path)
            elif file_type == 'presentations':
                return self.convert_pptx_to_pdf(input_path, output_path)
            elif file_type == 'pdf':
                # Copy existing PDF
                import shutil
                shutil.copy2(input_path, output_path)
                return True
            
            return False
        except Exception as e:
            print(f"Error converting {input_path}: {e}")
            return False


def generate_unique_filename(output_dir, base_name, extension=".pdf"):
    """Generate a unique filename to avoid conflicts in flat structure mode"""
    output_path = output_dir / f"{base_name}{extension}"
    counter = 1
    
    while output_path.exists():
        output_path = output_dir / f"{base_name}_{counter}{extension}"
        counter += 1
    
    return output_path


def combine_files_to_single_pdf(files, output_path, directory_path):
    """Combine all files into a single PDF with file titles"""
    try:
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.units import inch
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from PIL import Image
        import io
        
        # Create PDF document
        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Title style
        title_style = styles['Heading1']
        title_style.fontSize = 16
        title_style.spaceAfter = 12
        
        # Subtitle style for file names
        subtitle_style = styles['Heading2']
        subtitle_style.fontSize = 14
        subtitle_style.spaceAfter = 8
        
        # Normal style
        normal_style = styles['Normal']
        normal_style.fontSize = 10
        normal_style.spaceAfter = 6
        
        print(f"📄 Creating combined PDF: {output_path.name}")
        
        for i, file_path in enumerate(files, 1):
            try:
                relative_path = file_path.relative_to(directory_path)
                file_type = PDFConverter().get_file_type(file_path)
                
                print(f"[{i:3d}/{len(files)}] Adding: {relative_path}")
                
                # Add file title
                story.append(Paragraph(f"📄 {relative_path}", title_style))
                story.append(Spacer(1, 6))
                
                if file_type == 'images':
                    # For images, add a note and try to include the image
                    story.append(Paragraph(f"<i>Image file: {file_path.name}</i>", normal_style))
                    story.append(Paragraph(f"<i>Size: {file_path.stat().st_size} bytes</i>", normal_style))
                    
                    # Try to add the actual image
                    try:
                        with Image.open(file_path) as img:
                            # Convert to RGB if necessary
                            if img.mode != 'RGB':
                                img = img.convert('RGB')
                            
                            # Resize image to fit page width
                            max_width = A4[0] - 2*inch
                            max_height = A4[1] - 2*inch
                            
                            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
                            
                            # Save to bytes
                            img_bytes = io.BytesIO()
                            img.save(img_bytes, format='JPEG', quality=85)
                            img_bytes.seek(0)
                            
                            # Add image to PDF
                            from reportlab.platypus import Image as RLImage
                            story.append(RLImage(img_bytes, width=img.width, height=img.height))
                            story.append(Spacer(1, 12))
                    except Exception as e:
                        story.append(Paragraph(f"<i>Could not embed image: {str(e)}</i>", normal_style))
                
                elif file_type == 'documents':
                    # For text documents, add the content
                    try:
                        if file_path.suffix.lower() == '.txt':
                            with open(file_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                        elif file_path.suffix.lower() == '.md':
                            with open(file_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                        elif file_path.suffix.lower() == '.docx':
                            from docx import Document
                            doc_doc = Document(file_path)
                            content = '\n'.join([paragraph.text for paragraph in doc_doc.paragraphs])
                        else:
                            content = f"Document file: {file_path.name}"
                        
                        # Split content into paragraphs and add to PDF
                        lines = content.split('\n')
                        for line in lines:
                            if line.strip():
                                # Escape HTML special characters
                                escaped_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                story.append(Paragraph(escaped_line, normal_style))
                            else:
                                story.append(Spacer(1, 6))
                                
                    except Exception as e:
                        story.append(Paragraph(f"<i>Could not read document: {str(e)}</i>", normal_style))
                
                elif file_type == 'spreadsheets':
                    # For spreadsheets, add basic info
                    story.append(Paragraph(f"<i>Spreadsheet file: {file_path.name}</i>", normal_style))
                    story.append(Paragraph(f"<i>Size: {file_path.stat().st_size} bytes</i>", normal_style))
                    story.append(Paragraph(f"<i>Note: Spreadsheet content not fully displayed in combined PDF</i>", normal_style))
                
                elif file_type == 'presentations':
                    # For presentations, add basic info
                    story.append(Paragraph(f"<i>Presentation file: {file_path.name}</i>", normal_style))
                    story.append(Paragraph(f"<i>Size: {file_path.stat().st_size} bytes</i>", normal_style))
                    story.append(Paragraph(f"<i>Note: Presentation content not fully displayed in combined PDF</i>", normal_style))
                
                elif file_type == 'pdf':
                    # For existing PDFs, add basic info
                    story.append(Paragraph(f"<i>Existing PDF file: {file_path.name}</i>", normal_style))
                    story.append(Paragraph(f"<i>Size: {file_path.stat().st_size} bytes</i>", normal_style))
                    story.append(Paragraph(f"<i>Note: PDF content not embedded in combined PDF</i>", normal_style))
                
                # Add separator between files
                story.append(Spacer(1, 12))
                story.append(Paragraph("─" * 50, normal_style))
                story.append(Spacer(1, 12))
                
                # Add page break every 5 files to avoid too long pages
                if i % 5 == 0 and i < len(files):
                    story.append(PageBreak())
                
            except Exception as e:
                story.append(Paragraph(f"<i>Error processing {file_path.name}: {str(e)}</i>", normal_style))
                story.append(Spacer(1, 12))
        
        # Build the PDF
        doc.build(story)
        return True
        
    except Exception as e:
        print(f"Error creating combined PDF: {e}")
        return False


def convert_directory(directory, maintain_structure=True, output_dir=None, verbose=False, combine=False):
    """Convert all files in directory and subdirectories"""
    converter = PDFConverter()
    
    directory_path = Path(directory).resolve()
    
    if not directory_path.exists():
        print(f"❌ Error: Directory '{directory}' does not exist")
        return False
    
    if not directory_path.is_dir():
        print(f"❌ Error: '{directory}' is not a directory")
        return False
    
    # Determine output directory
    if output_dir:
        pdf_output_dir = Path(output_dir).resolve()
    else:
        pdf_output_dir = directory_path / "pdf"
    
    # Create PDF output directory
    pdf_output_dir.mkdir(exist_ok=True)
    
    structure_mode = "directory structure" if maintain_structure else "flat list"
    if combine:
        structure_mode = "combined single PDF"
    
    print(f"🔄 PDF Converter - Command Line")
    print(f"📁 Source directory: {directory_path}")
    print(f"💾 Output directory: {pdf_output_dir}")
    print(f"🏗️  Structure mode: {structure_mode}")
    print()
    
    # Get all files to convert
    all_files = []
    for root, dirs, files in os.walk(directory_path):
        for file in files:
            file_path = Path(root) / file
            if converter.get_file_type(file_path):
                all_files.append(file_path)
    
    if not all_files:
        print("❌ No supported files found in the directory")
        print("\nSupported formats:")
        print("• Images: JPG, PNG, GIF, BMP, TIFF, WebP")
        print("• Documents: DOCX, TXT, Markdown")
        print("• Spreadsheets: XLSX, XLS")
        print("• Presentations: PPTX")
        print("• PDFs: Existing PDFs will be copied")
        return False
    
    print(f"📋 Found {len(all_files)} files to convert:")
    if verbose:
        for file_path in all_files:
            relative_path = file_path.relative_to(directory_path)
            print(f"   • {relative_path}")
    else:
        print(f"   Use --verbose to see file list")
    print()
    
    # Handle combine mode
    if combine:
        output_filename = f"{directory_path.name}_combined.pdf"
        combined_output_path = pdf_output_dir / output_filename
        
        print(f"📄 Combining all files into single PDF...")
        print(f"📄 Output file: {combined_output_path}")
        print()
        
        if combine_files_to_single_pdf(all_files, combined_output_path, directory_path):
            print("✅ Combined PDF created successfully!")
            print(f"📁 Combined PDF saved to: {combined_output_path}")
            return True
        else:
            print("❌ Failed to create combined PDF")
            return False
    
    successful_conversions = 0
    failed_conversions = 0
    
    for i, file_path in enumerate(all_files, 1):
        try:
            # Create relative path structure in PDF directory
            relative_path = file_path.relative_to(directory_path)
            pdf_filename = relative_path.stem + ".pdf"
            
            if maintain_structure:
                # Maintain directory structure
                output_subdir = pdf_output_dir / relative_path.parent
                output_subdir.mkdir(parents=True, exist_ok=True)
                output_path = output_subdir / pdf_filename
            else:
                # Flat structure - all files in root pdf directory
                base_name = relative_path.stem
                # If file is in subdirectory, include parent directory name to make it unique
                if relative_path.parent != Path('.'):
                    parent_name = str(relative_path.parent).replace('/', '_').replace('\\', '_')
                    base_name = f"{parent_name}_{base_name}"
                
                output_path = generate_unique_filename(pdf_output_dir, base_name)
            
            print(f"[{i:3d}/{len(all_files)}] Converting: {relative_path}", end=" ")
            
            # Convert file
            if converter.convert_file(str(file_path), str(output_path)):
                successful_conversions += 1
                if maintain_structure:
                    print(f"✅")
                    if verbose:
                        print(f"           → {output_path.relative_to(pdf_output_dir)}")
                else:
                    print(f"✅ → {output_path.name}")
            else:
                failed_conversions += 1
                print(f"❌")
                
        except Exception as e:
            failed_conversions += 1
            print(f"❌ Error: {str(e)}")
    
    # Final summary
    print()
    print("📊 Conversion Summary:")
    print(f"   ✅ Successful: {successful_conversions}")
    print(f"   ❌ Failed: {failed_conversions}")
    print(f"   📁 PDFs saved to: {pdf_output_dir}")
    print(f"   🏗️  Structure: {structure_mode}")
    
    if successful_conversions > 0:
        print(f"\n🎉 Conversion completed successfully!")
        return True
    else:
        print(f"\n💔 No files were converted successfully")
        return False


def main():
    """Main function for command line interface"""
    parser = argparse.ArgumentParser(
        description="Convert files in a directory to PDF format",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s /path/to/documents
  %(prog)s /path/to/documents --flat
  %(prog)s /path/to/documents --output /path/to/pdfs
  %(prog)s /path/to/documents --flat --verbose
  %(prog)s /path/to/documents --combine
  %(prog)s . --output ../converted_pdfs

Supported file formats:
  • Images: JPG, PNG, GIF, BMP, TIFF, WebP
  • Documents: DOCX, TXT, Markdown
  • Spreadsheets: XLSX, XLS  
  • Presentations: PPTX
  • PDFs: Existing PDFs will be copied
        """
    )
    
    parser.add_argument(
        "directory",
        help="Directory containing files to convert"
    )
    
    parser.add_argument(
        "--flat",
        action="store_true",
        help="Create flat structure (all PDFs in one folder) instead of maintaining directory structure"
    )
    
    parser.add_argument(
        "--output", "-o",
        help="Output directory for PDFs (default: creates 'pdf' folder in source directory)"
    )
    
    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Show detailed output including file list and paths"
    )
    
    parser.add_argument(
        "--combine", "-c",
        action="store_true",
        help="Combine all files into a single PDF with file titles"
    )
    
    parser.add_argument(
        "--version",
        action="version",
        version="PDF Converter CLI 1.0"
    )
    
    args = parser.parse_args()
    
    try:
        success = convert_directory(
            directory=args.directory,
            maintain_structure=not args.flat,
            output_dir=args.output,
            verbose=args.verbose,
            combine=args.combine
        )
        
        sys.exit(0 if success else 1)
        
    except KeyboardInterrupt:
        print("\n\n⚠️  Conversion interrupted by user")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ Unexpected error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main() 